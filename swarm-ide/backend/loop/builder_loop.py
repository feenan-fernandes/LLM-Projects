"""
builder_loop.py
8-iteration state machine implementing the Section 5 Orchestrator action loop.
Handles all action tags, enforces the install_skill human-approval gate,
delegates self_heal to the Tester agent, and enforces the iter-6 partial rule.
"""
import os
import json
import re

from backend.agents.orchestrator import call_orchestrator
from backend.agents.tester import diagnose_and_fix
from backend.agents.skill_installer import install_skill, load_skill_context
from backend.loop.action_parser import parse_action, extract_thought
from backend.loop.sandbox import execute_command_safely, SandboxViolationError
from backend.governance.logger import log_action

MAX_ITERATIONS = 10
PARTIAL_THRESHOLD = 6   # if no tests passing by this iter â†’ flag partial


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _observation_prefix(action_type, iteration):
    return f"[iter {iteration}] {action_type.upper()} result"


import py_compile
def validate_file(path):
    if not path.endswith('.py'): return None
    try:
        py_compile.compile(path, doraise=True)
        return None
    except Exception as e:
        return "Syntax Error in file: " + str(e)

def _human_approval_pending(action):
    """Returns an ArtifactCard-style dict the caller must surface to the user."""
    return {
        "requires_approval": True,
        "action_type": action["type"],
        "args": action["args"],
        "message": (
            "Human approval required before this action executes. "
            "Review in ArtifactCard and click Approve to proceed."
        ),
    }


# ---------------------------------------------------------------------------
# Main loop
# ---------------------------------------------------------------------------


def _truncate_output(text: str, max_chars: int = 3000) -> str:
    if not text: return ""
    if len(text) <= max_chars: return text
    half = max_chars // 2
    return text[:half] + "\n\n... [OUTPUT TRUNCATED to protect context window. Use grep to search] ...\n\n" + text[-half:]

def run_builder_loop(
    task_description,
    history=None,
    task_id="task_default",
    mock_responses=None,
    # Callbacks provided by the Flask route / test harness:
    approval_callback=None,   # callable(pending_dict) -> bool  (True = approved)
    stream_callback=None,     # callable(event_dict) -> None    (for SSE)
    skill_context_override=None,  # pre-loaded SKILL.md text injected at start
    model="deepseek-r1:7b",
    system_prompt=None,
    abort_event=None
):
    """
    Executes up to MAX_ITERATIONS of the Section 5 action loop.

    Returns:
        (success: bool, iterations: int, summary: str)
    """
    conversation = f"USER REQUEST: {task_description}\n"
    if skill_context_override:
        conversation += f"\nSKILL CONTEXT:\n{skill_context_override}\n"
    conversation += "\nWhat is your first action?"

    tests_passed = False
    final_summary = ""
    consecutive_invalid = 0
    total_tokens = 0
    max_tps = 0.0

    def _emit(event):
        if stream_callback:
            stream_callback(event)

    for iteration in range(1, MAX_ITERATIONS + 1):
        if iteration > 1 and iteration % 6 == 0:
            _emit({"type": "system", "msg": f"Step {iteration}: Context window large. Running Auto-Compaction..."})
            try:
                from backend.agents.compactor import compact_trajectory, apply_compaction
                mem_block = compact_trajectory(conversation, model=model)
                conversation = apply_compaction(conversation, mem_block, keep_last_n_turns=2)
                _emit({"type": "action", "iteration": iteration, "action": "memory_compaction", "result": "Context window limits approached.\nTrajectory successfully compressed into <working_memory> block."})
            except Exception as e:
                _emit({"type": "system", "msg": f"Auto-Compaction failed: {e}"})

        if abort_event and abort_event.is_set():
            _emit({"type": "action", "iteration": iteration, "action": "abort", "result": "Task aborted by user."})
            return False, iteration, "Aborted by user", total_tokens, max_tps

        i = iteration - 1
        _emit({"type": "system", "msg": f"Step {iteration}: Agent is thinking..."})

        # --- Call Orchestrator (or use mock) ---
        mock = mock_responses[i] if mock_responses and i < len(mock_responses) else None
        response_text, metrics = call_orchestrator(conversation, model=model, mock_response=mock, system_prompt=system_prompt)
        
        if not response_text.strip() and not mock:
            summary = "Error: Local LLM returned an empty response. Ensure Ollama is running and the model is downloaded."
            _emit({"type": "action", "iteration": iteration, "action": "abort", "result": summary})
            return False, iteration, summary, total_tokens, max_tps

        c_tokens = metrics.get('completion_tokens', 0)
        e_dur = metrics.get('eval_duration', 1) / 1e9
        total_tokens += c_tokens
        if e_dur > 0:
            tps = c_tokens / e_dur
            if tps > max_tps:
                max_tps = tps

        thought = extract_thought(response_text)
        action = parse_action(response_text)

        _emit({
            "type": "thought",
            "iteration": iteration,
            "thought": thought,
        })

        # --- No valid action ---
        if not action:
            consecutive_invalid += 1
            if consecutive_invalid >= 3:
                summary = "Failed: 3 consecutive invalid XML actions."
                _emit({"type": "action", "iteration": iteration, "action": "invalid_xml", "result": summary})
                return False, iteration, summary, total_tokens, max_tps

            # Partial check at iter 6
            if iteration >= PARTIAL_THRESHOLD and not tests_passed:
                summary = (
                    "Reached iteration 6 without passing tests and no valid action tag. "
                    "Flagging as partial."
                )
                log_action(task_id, iteration, "partial_flag", summary, 0, 0)
                _emit({"type": "action", "iteration": iteration, "action": "partial_flag", "result": summary})
                return False, iteration, summary, total_tokens, max_tps

            obs = (
            "ERROR: No valid XML action tag found. "
            "You MUST output exactly ONE valid action tag. "
            "Do NOT just converse. Use <write_file>, <execute_bash>, <patch_file>, <create_ppt> or <finish>.\n"
            "If the user asks for a PowerPoint or presentation, DO NOT write a python script. ALWAYS use the <create_ppt> tool directly.\n"
            "If the user asks for a wireframe or UI, use <write_file> to generate an .html file with Tailwind CSS. DO NOT try to write a python UI script.\n"
            "Example PPTX Tool Call:\n<create_ppt>\n<path>presentation.pptx</path>\n<content>Title Slide\nSubtitle\n---SLIDE---\nSlide 2\nBullet point</content>\n</create_ppt>\n"
            "Example:\n<execute_bash>\n  <command>ls -la</command>\n</execute_bash>"
        )
            log_action(task_id, iteration, "invalid_xml", obs, metrics["completion_tokens"], metrics["eval_duration"] // 1_000_000)
            # Do not append hallucinated garbage to context, it causes 7B models to spiral
            conversation += f"\n\nSYSTEM OBSERVATION:\n{obs}\nDo not repeat the invalid output. Next action?"
            _emit({"type": "action", "iteration": iteration, "action": "invalid_xml", "result": obs})
            continue

        consecutive_invalid = 0
        action_type = action["type"]
        
        if action_type == "plan" and "<plan>" in conversation:
            observation = "ERROR: You already submitted a plan in a previous step! You MUST NOT submit another <plan>. Use <spawn_worker>, <write_file>, <execute_bash>, etc."
            # Do not append hallucinated plan repetition to context
            conversation += f"\n\nSYSTEM OBSERVATION:\n{observation}\nDo not repeat the plan. Take action."
            _emit({"type": "action", "iteration": iteration, "action": "invalid_plan_loop", "result": observation})
            continue
        args = action["args"]
        observation = ""
        is_finished = False

        # ---------------------------------------------------------------
        try:
            # 1. PLAN â€” surface to UI, no side-effects
            if action_type == "plan":
                observation = (
                    f"Plan acknowledged.\n"
                    f"Goal: {args['goal']}\n"
                    f"Steps: {args['steps']}\n"
                    f"Acceptance: {args['acceptance_criteria']}\n\n"
                    f"IMPORTANT: Plan is stored in memory. DO NOT output <plan> again. You must now take your first action using one of your available tools (e.g. <spawn_worker>, <execute_bash>, <write_file>, etc)."
                )

            # 2. SEARCH GITHUB SKILLS â€” read-only, no approval needed
            elif action_type == "search_github_skills":
                from backend.agents.skill_scout import search_skills
                results = search_skills(
                    args["query"],
                    max_candidates=5,
                )
                if results:
                    names = [r["full_name"] for r in results]
                    observation = f"Found {len(results)} candidate skill(s): {names}"
                else:
                    observation = "No candidate skills found matching query."
                _emit({"type": "skill_discovery", "candidates": results if results else []})

            # 3. EVALUATE SKILL â€” Librarian grades one repo
            elif action_type == "evaluate_skill":
                from backend.agents.librarian import grade_skill_relevance
                repo_url = args.get("repo_url", "")
                # In real flow we'd fetch SKILL.md; in loop we use repo description
                approved = grade_skill_relevance(task_description, repo_url)
                observation = f"Librarian evaluation: {'APPROVED' if approved else 'REJECTED'} â€” {repo_url}"

            # 4. INSTALL SKILL â€” REQUIRES HUMAN APPROVAL
            elif action_type == "install_skill":
                pending = _human_approval_pending(action)
                _emit({"type": "approval_required", "pending": pending})

                approved = False
                if approval_callback:
                    approved = approval_callback(pending)

                if approved:
                    # Build a minimal candidate dict from args
                    slug = args["repo_url"].rstrip("/").split("/")[-1]
                    candidate = {
                        "slug": slug,
                        "name": slug,
                        "full_name": args["repo_url"].replace("https://github.com/", ""),
                        "url": args["repo_url"],
                        "stars": 0,
                        "score": 0.0,
                        "skill_md": f"# {slug}\nInstalled from {args['repo_url']}",
                    }
                    path = install_skill(candidate, task_id=task_id, log_fn=log_action)
                    observation = f"Skill installed to {path}."
                else:
                    observation = "Skill install REJECTED by human reviewer. Skipping."
                    log_action(task_id, iteration, "install_rejected", observation, 0, 0)

            # 5. SELECT SKILL
            elif action_type == "select_skill":
                slug = args.get("skill_name", "")
                context = load_skill_context(slug)
                if context:
                    conversation += f"\n\nSELECTED SKILL CONTEXT ({slug}):\n{context}\n"
                    observation = f"Skill '{slug}' selected and context injected."
                else:
                    observation = f"Skill '{slug}' not found in registry. Has it been installed?"

            # 5.5 PATCH FILE
            elif action_type == "patch_file":
                path = args.get("path", "")
                diff_str = args.get("diff", "")
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                full_path = os.path.abspath(os.path.join(workspace_dir, path))
                
                path_safe = False
                try:
                    path_safe = os.path.commonpath([workspace_dir, full_path]) == workspace_dir
                except ValueError:
                    path_safe = False
                
                if not path_safe:
                    observation = f"Sandbox violation: Cannot write outside workspace directory."
                else:
                    try:
                        import patch as patch_lib
                        pset = patch_lib.fromstring(diff_str.encode('utf-8'))
                        success = pset.apply(root=os.path.dirname(full_path))
                        if success:
                            observation = f"Patch applied successfully to {path}"
                            err = validate_file(full_path)
                            if err: observation += "\n" + err
                        else:
                            observation = f"Failed to apply patch to {path}. Context mismatch. Try <write_file> instead."
                    except Exception as e:
                        observation = f"Error applying patch to {path}: {e}"

            # 6. WRITE FILE

            elif action_type == "replace_block":
                path = args.get("path", "")
                search = args.get("search", "")
                replace = args.get("replace", "")
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                full_path = os.path.abspath(os.path.join(workspace_dir, path))
                
                try:
                    path_safe = os.path.commonpath([workspace_dir, full_path]) == workspace_dir
                except ValueError:
                    path_safe = False
                
                if not path_safe:
                    observation = f"Sandbox violation: Cannot write outside workspace directory."
                else:
                    try:
                        with open(full_path, "r", encoding="utf-8") as f:
                            content = f.read()
                        if search not in content:
                            observation = f"Failed to patch: The exact search block was not found in {path}. Make sure whitespace matches exactly."
                        else:
                            content = content.replace(search, replace, 1)
                            with open(full_path, "w", encoding="utf-8") as f:
                                f.write(content)
                            observation = f"Successfully replaced block in {path}."
                            err = validate_file(full_path)
                            if err: observation += "\n" + err
                    except Exception as e:
                        observation = f"Error reading/writing {path}: {e}"
            elif action_type == "create_ppt":
                path = args.get("path", "presentation.pptx")
                content = args.get("content", "")
                
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                full_path = os.path.abspath(os.path.join(workspace_dir, path))
                
                try:
                    from pptx import Presentation
                    prs = Presentation()
                    # A very basic parser to split content by slides
                    slides_text = content.split('---SLIDE---')
                    for slide_txt in slides_text:
                        if not slide_txt.strip(): continue
                        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title & Content layout
                        lines = [l.strip() for l in slide_txt.strip().split('\n') if l.strip()]
                        if lines:
                            slide.shapes.title.text = lines[0]
                            if len(lines) > 1:
                                slide.placeholders[1].text = '\n'.join(lines[1:])
                    
                    os.makedirs(os.path.dirname(full_path), exist_ok=True)
                    prs.save(full_path)
                    observation = f"PowerPoint created natively: {path}"
                except Exception as e:
                    observation = f"Error generating PPTX: {e}"

            elif action_type == "write_file":
                path = args.get("path", "")
                content = args.get("content", "")
                
                # Enforce workspace boundary
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                full_path = os.path.abspath(os.path.join(workspace_dir, path))
                path_safe = False
                try:
                    path_safe = os.path.commonpath([workspace_dir, full_path]) == workspace_dir
                except ValueError:
                    path_safe = False
                
                if not path_safe:
                    observation = f"Sandbox violation: Cannot write outside workspace directory."
                else:
                    os.makedirs(os.path.dirname(full_path), exist_ok=True)
                    with open(full_path, "w", encoding="utf-8") as f:
                        f.write(content)
                    observation = f"File written: {path}"
                    err = validate_file(full_path)
                    if err: observation += "\n" + err

            elif action_type == "create_skill":
                name = args.get("name", "")
                code = args.get("code", "")
                if not name.endswith(".py"): name += ".py"
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                skills_dir = os.path.join(workspace_dir, "skills")
                os.makedirs(skills_dir, exist_ok=True)
                
                full_path = os.path.join(skills_dir, os.path.basename(name))
                with open(full_path, "w", encoding="utf-8") as f:
                    f.write(code)
                observation = f"Skill created and persisted permanently at workspace/skills/{os.path.basename(name)}"

            elif action_type == "use_skill":
                name = args.get("name", "")
                skill_args = args.get("args", "")
                if not name.endswith(".py"): name += ".py"
                workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                full_path = os.path.join(workspace_dir, "skills", os.path.basename(name))
                
                if not os.path.exists(full_path):
                    observation = f"Skill '{name}' does not exist in workspace/skills/. You must <create_skill> first."
                else:
                    import shlex
                    safe_args = shlex.split(skill_args) if skill_args else []
                    cmd = f'python "{full_path}" ' + ' '.join(f'"{a}"' for a in safe_args)
                    res = execute_command_safely(cmd, task_id=task_id, iteration=iteration)
                    observation = f"Skill Execution STDOUT:\n{_truncate_output(res['stdout'])}\nSTDERR:\n{_truncate_output(res['stderr'])}\nExit: {res['code']}"

            # 7. EXECUTE BASH
            elif action_type == "search_web":
                query = args.get("query", "")
                try:
                    from duckduckgo_search import DDGS
                    results = DDGS().text(query, max_results=3)
                    if results:
                        observation = "Search Results:\n"
                        for r in results:
                            observation += f"- [{{r.get('title')}}]({{r.get('href')}})\n{{r.get('body')}}\n\n"
                    else:
                        observation = "No results found for query."
                except Exception as e:
                    observation = f"Web search failed: {e}"

            elif action_type == "mcp_call":
                try:
                    from backend.loop.mcp_bridge import get_mcp_bridge
                    workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..', 'workspace'))
                    bridge = get_mcp_bridge(workspace_dir)
                    server = args.get("server", "unknown")
                    tool = args.get("tool", "unknown")
                    arguments = args.get("arguments", {})
                    observation = bridge.call_tool_sync(server, tool, arguments)
                except Exception as e:
                    observation = f"MCP call failed: {e}"

            elif action_type == "spawn_worker":
                tasks = args.get("tasks", [])
                import concurrent.futures
                
                def _run_subtask(idx, t):
                    _emit({"type": "system", "msg": f"Sub-Agent {idx+1} spawned: {t[:30]}..."})
                    prompt = f"You are a fast Swarm Sub-Agent. Your isolated task is: {t}\nReturn a thorough but concise answer."
                    try:
                        res, _ = call_orchestrator(prompt, model=model)
                        return f"--- Sub-Agent {idx+1} Result ---\n{res}\n"
                    except Exception as e:
                        return f"--- Sub-Agent {idx+1} Failed ---\n{e}\n"
                
                with concurrent.futures.ThreadPoolExecutor(max_workers=min(4, len(tasks))) as executor:
                    futures = []
                    for idx, t in enumerate(tasks):
                        futures.append(executor.submit(_run_subtask, idx, t))
                    results = [f.result() for f in futures]
                
                observation = "Sub-agents completed their parallel tasks:\n\n" + "\n".join(results)
                _emit({"type": "system", "msg": f"All {len(tasks)} sub-agents completed."})

            elif action_type == "execute_bash":
                cmd = args.get("command", "")
                res = execute_command_safely(cmd, task_id=task_id, iteration=iteration)
                observation = f"STDOUT:\n{_truncate_output(res['stdout'])}\nSTDERR:\n{_truncate_output(res['stderr'])}\nExit: {res['code']}"

            # 8. RUN TEST (Phase 5 InspectCoder augmentation)
            elif action_type == "run_test":
                cmd = args.get("command", "")
                expect = args.get("expect", "")
                
                # Phase 5: Automatically inject --showlocals if using pytest to capture runtime variables
                if "pytest" in cmd and "--showlocals" not in cmd and "-l" not in cmd:
                    cmd += " --showlocals"
                    
                res = execute_command_safely(cmd, task_id=task_id, iteration=iteration)
                passed = res["code"] == 0
                tests_passed = tests_passed or passed
                observation = (
                    f"Test {'PASSED' if passed else 'FAILED'}.\n"
                    f"Expected: {expect}\n"
                    f"STDOUT:\n{res['stdout']}\nSTDERR:\n{res['stderr']}"
                )
                _emit({"type": "test_result", "passed": passed, "output": observation})

                # Iteration 6 partial rule
                if iteration >= PARTIAL_THRESHOLD and not tests_passed:
                    summary = (
                        "Reached iteration 6 without passing tests. "
                        "Narrowing scope and flagging partial."
                    )
                    log_action(task_id, iteration, "partial_flag", summary, 0, 0)
                    _emit({"type": "action", "iteration": iteration, "action": "partial_flag", "result": summary})
                    return False, iteration, summary, total_tokens, max_tps

            # 9. SELF HEAL â€” delegate to Tester
            elif action_type == "self_heal":
                error_summary = args.get("error_summary", "")
                hypothesis = args.get("hypothesis", "")
                # Tester needs the last written file; we pass the hypothesis as context
                fixed = diagnose_and_fix(
                    original_code=hypothesis,
                    error_trace=error_summary,
                )
                observation = f"Tester self-heal response:\n{fixed}"


            # 10. FINISH
            elif action_type == "finish":
                status = args.get("status", "success")
                summary = args.get("summary", "")
                final_summary = summary
                is_finished = True
                observation = f"Finished [{status}]: {summary}"

        except SandboxViolationError as e:
            observation = str(e)
        except Exception as e:
            observation = f"System error in {action_type}: {e}"

        # --- Governance log ---
        log_action(
            task_id, iteration, action_type,
            observation[:2000],          # cap content to keep db lean
            metrics.get("completion_tokens", 0),
            metrics.get("eval_duration", 0) // 1_000_000,
        )

        _emit({
            "type": "action",
            "iteration": iteration,
            "action": action_type,
            "result": observation,
            "metrics": metrics,
        })

        if is_finished:
            success = action["args"].get("status", "success") == "success"
            return success, iteration, final_summary, total_tokens, max_tps

        # Append observation to conversation
        response_text = re.sub(r'<think>.*?</think>', '', response_text, flags=re.DOTALL).strip()
        conversation += (
            f"\n\nASSISTANT:\n{response_text}\n\n"
            f"<observation>\n{observation}\n</observation>\n\n"
            "Next action?"
        )

    # Iteration budget exhausted
    log_action(task_id, MAX_ITERATIONS, "timeout", "Iteration budget (8) exhausted.", 0, 0)
    return False, MAX_ITERATIONS, "Failed: iteration budget exhausted.", total_tokens, max_tps

